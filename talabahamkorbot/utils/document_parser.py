
import os
import io
import logging
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)

def extract_text_from_stream(file_stream, file_ext: str) -> str:
    """
    Turli fayl formatlaridan (PDF, DOCX, PPTX, TXT) matnni ajratib oladi.
    file_stream: file-like object (BytesIO, opened file, etc.)
    """
    full_text = ""
    file_ext = file_ext.lower().replace(".", "")

    try:
        # Move request to beginning of stream if needed
        if hasattr(file_stream, 'seek'):
            file_stream.seek(0)
            
        if file_ext == "pdf":
            reader = PdfReader(file_stream)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    full_text += text + "\n"
        
        elif file_ext == "docx":
            doc = Document(file_stream)
            for para in doc.paragraphs:
                full_text += para.text + "\n"
        
        elif file_ext == "pptx":
            prs = Presentation(file_stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text += shape.text + "\n"
        
        elif file_ext == "txt":
            # Try decoding utf-8
            content = file_stream.read()
            if isinstance(content, bytes):
                full_text = content.decode("utf-8")
            else:
                full_text = str(content)

        else:
            return "Kechirasiz, bu fayl formati hozircha qo'llab-quvvatlanmaydi."

    except Exception as e:
        logger.error(f"Faylni o'qishda xatolik: {e}")
        return "Faylni o'qishda xatolik yuz berdi. Fayl shikastlangan bo'lishi mumkin."

    return full_text.strip()

def extract_text_from_file(file_path: str, file_ext: str) -> str:
    """
    Legacy wrapper for file paths.
    """
    with open(file_path, "rb") as f:
        return extract_text_from_stream(f, file_ext)
